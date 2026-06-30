import streamlit as st
import openpyxl
import docx
import json
from pdf2docx import Converter
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Cm, Mm
from docx.text.paragraph import Paragraph
from pptx import Presentation
from pdf2docx import Converter # THƯ VIỆN MỚI GIỮ ĐỊNH DẠNG PDF
import google.generativeai as genai
from google.generativeai.types import HarmCategory, HarmBlockThreshold
import time
import io
import os
import re
from PIL import Image
from streamlit_paste_button import paste_image_button
import pandas as pd 

# ==========================================
# BỘ LỌC AN TOÀN - TẮT 100% KIỂM DUYỆT
# ==========================================
SAFETY_SETTINGS = {
    HarmCategory.HARM_CATEGORY_HARASSMENT: HarmBlockThreshold.BLOCK_NONE,
    HarmCategory.HARM_CATEGORY_HATE_SPEECH: HarmBlockThreshold.BLOCK_NONE,
    HarmCategory.HARM_CATEGORY_SEXUALLY_EXPLICIT: HarmBlockThreshold.BLOCK_NONE,
    HarmCategory.HARM_CATEGORY_DANGEROUS_CONTENT: HarmBlockThreshold.BLOCK_NONE,
}

# ==========================================
# DỌN DẸP RÁC HỆ THỐNG KHI KHỞI ĐỘNG
# ==========================================
temp_files = [f for f in os.listdir('.') if f.startswith('temp_')]
for f in temp_files:
    try: os.remove(f)
    except: pass

# ==========================================
# CẤU HÌNH GIAO DIỆN
# ==========================================
st.set_page_config(page_title="AI HMI Translator Pro", layout="wide")
st.title("🌐 Ứng Dụng Dịch Tài Liệu Kỹ Thuật Đa Năng")
st.markdown("Hỗ trợ: **Word, Excel, PowerPoint, PDF** | Thuật toán Nhồi Dữ Liệu Động | Dịch Ảnh Định Vị.")

# ==========================================
# SIDEBAR & MENU CÔNG CỤ
# ==========================================
with st.sidebar:
    st.header("🔑 Cấu Hình Hệ Thống")
    try:
        api_key = st.secrets["GEMINI_API_KEY"]
        genai.configure(api_key=api_key)
    except:
        api_key = st.text_input("Nhập Google Gemini API Key:", type="password")
        if api_key: genai.configure(api_key=api_key)
            
    selected_model_name = "gemini-3.1-flash-lite"
    if api_key:
        try:
            available_models = [m.name.replace('models/', '') for m in genai.list_models() if 'generateContent' in m.supported_generation_methods]
            if available_models:
                st.success("✅ Đã kết nối API thành công!")
                default_idx = 0
                for i, m_name in enumerate(available_models):
                    if m_name.strip().lower() == "gemini-3.1-flash-lite":
                        default_idx = i
                        break
                selected_model_name = st.selectbox("🤖 Chọn Model AI:", available_models, index=default_idx)
        except Exception:
            selected_model_name = st.selectbox("🤖 Chọn Model AI (Thủ công):", ["gemini-3.1-flash-lite", "gemini-1.5-flash"], index=0)

    st.markdown("---")
    st.header("🛠 Tùy Chọn Chạy File")
    
    smart_resume = st.toggle("🧠 Dịch Nối Tiếp", value=True, help="Quét tia X cực nhanh: Tự động bỏ qua các câu đã là tiếng Việt/tiếng Anh. Khuyên dùng khi kéo thả file nháp vào để dịch tiếp.")
    skip_shapes = st.toggle("⏭️ Bỏ qua chữ trong Lưu đồ", value=False, help="Bật lên nếu file Word của bạn có nhiều sơ đồ dễ bị vỡ khung. Máy sẽ chừa lại không dịch các khối này.") 
    
    st.markdown("---")
    st.header("🎯 Giới Hạn Vùng Dịch")
    dich_theo_vung = st.toggle("📍 Kích hoạt chọn vùng dịch", value=False, help="Chỉ dịch từ đoạn X đến đoạn Y. Phải dùng kết hợp với nút 'SOI VỊ TRÍ' ở màn hình chính để biết chính xác số đoạn.")
    
    start_idx, end_idx = 1, 0
    if dich_theo_vung:
        col1, col2 = st.columns(2)
        with col1: start_idx = st.number_input("Từ mục số:", min_value=1, value=1)
        with col2: end_idx = st.number_input("Đến mục số:", min_value=1, value=100)

    with st.expander("📖 Cẩm Nang Hướng Dẫn Tính Năng"):
        st.markdown("""
        **1. 🧠 Dịch Nối Tiếp (Smart Resume):**
        - Dùng tia X quét file: Đoạn nào không có tiếng Trung sẽ bỏ qua luôn (0.001s).
        - Tuyệt chiêu: File rớt mạng đang dịch dở tải về, thả vào dịch tiếp sẽ không bị tính lại từ đầu.
        
        **2. ⏭️ Bỏ qua chữ Lưu đồ:**
        - Nếu bật: Máy sẽ không dịch chữ trong các khối hình vẽ (Shapes) để tránh làm vỡ khung.
        
        **3. 📍 Chọn vùng dịch & Quét X-Quang:**
        - Bạn được xem trước danh sách file được đánh số thứ tự từ 1 đến hết. 
        - Nhập số vào ô bên trên để máy CHỈ dịch đúng đoạn đó, tiết kiệm thời gian.
        
        **4. ⚡ Thuật toán Nhồi Động (Dynamic Chunking):**
        - Chạy ngầm phía sau. Tự động đo độ dài đoạn văn, gom câu ngắn, tách câu dài cho vừa đúng 1500 ký tự để ép tốc độ dịch lên tối đa mà không bị Google khóa.
        """)

# ==========================================
# CÁC HÀM XỬ LÝ LÕI & DYNAMIC CHUNKING (CHỐNG KẸT API)
# ==========================================
def handle_api_error(e, attempt):
    error_msg = str(e).lower()
    if "429" in error_msg or "quota" in error_msg or "exhausted" in error_msg:
        st.toast("⚠️ Quá tải API! Đang nghỉ đông 60s để phục hồi...", icon="⏳")
        time.sleep(60) 
    elif "safety" in error_msg or "block" in error_msg:
        st.toast("🛑 Lỗi Kiểm Duyệt: AI từ chối cụm từ nhạy cảm!", icon="🚫")
        time.sleep(2)
    else:
        st.toast(f"Lỗi API. Thử lại lần {attempt+1}...", icon="🔄")
        time.sleep(5)

def translate_text_core(text, model_name):
    prompt = f"Dịch đoạn văn bản kỹ thuật sau từ tiếng Trung sang tiếng Việt. KHÔNG giải thích.\nGốc:\n{text}\nDịch:"
    model = genai.GenerativeModel(model_name)
    for attempt in range(3):
        try:
            time.sleep(5) # ÉP TỐC ĐỘ 12 REQ/PHÚT
            res = model.generate_content(prompt, safety_settings=SAFETY_SETTINGS)
            if res.text: return res.text.replace("Bản dịch tiếng Việt:", "").strip()
        except Exception as e:
            handle_api_error(e, attempt)
    return text

def translate_batch_core(texts, model_name):
    if not texts: return []
    if len(texts) == 1: return [translate_text_core(texts[0], model_name)]
    
    # SỬ DỤNG JSON ĐỂ CHỐNG LỖI MẤT ĐOẠN CỦA AI
    json_input = json.dumps(texts, ensure_ascii=False)
    prompt = f"""Dịch mảng JSON chứa văn bản kỹ thuật sau sang tiếng Việt.
YÊU CẦU BẮT BUỘC:
1. Trả về ĐÚNG CẤU TRÚC mảng JSON.
2. Số lượng phần tử đầu ra PHẢI BẰNG ĐÚNG đầu vào ({len(texts)} phần tử).
3. KHÔNG có markdown, KHÔNG giải thích. Chỉ in ra mảng JSON hợp lệ.

Đầu vào:
{json_input}

Đầu ra:"""

    model = genai.GenerativeModel(model_name)
    for attempt in range(3):
        try:
            time.sleep(5) 
            response = model.generate_content(prompt, safety_settings=SAFETY_SETTINGS)
            if response.text:
                res_text = response.text.strip()
                # Lọc an toàn không dùng 3 dấu nháy ngược gây lỗi giao diện
                res_text = res_text.strip("` \n")
                if res_text.lower().startswith("json"):
                    res_text = res_text[4:].strip()
                
                try:
                    translated_array = json.loads(res_text)
                    if isinstance(translated_array, list) and len(translated_array) == len(texts):
                        return translated_array
                except Exception:
                    pass # Bỏ qua để thử lại nếu JSON hỏng
        except Exception as e:
            handle_api_error(e, attempt)
            
    # NẾU JSON VẪN THẤT BẠI: CHIA ĐÔI MẢNG THAY VÌ GỌI TỪNG CÂU (CHỐNG BÙNG NỔ REQUEST)
    st.toast("⚠️ Phân tách gói dịch để tránh lỗi API...", icon="🔧")
    mid = len(texts) // 2
    return translate_batch_core(texts[:mid], model_name) + translate_batch_core(texts[mid:], model_name)

def dynamic_chunking(items, max_chars=1500, max_items=25):
    batches = []
    current_batch = []
    current_len = 0
    for item in items:
        clean_text = item['text'].replace('|||', '---')
        item['text'] = clean_text
        l = len(clean_text)
        
        # Thêm giới hạn max_items=25 để tránh AI bị ngộp khi xử lý mảng JSON quá dài
        if (current_len + l > max_chars or len(current_batch) >= max_items) and current_batch:
            batches.append(current_batch)
            current_batch = [item]
            current_len = l
        else:
            current_batch.append(item)
            current_len += l
    if current_batch:
        batches.append(current_batch)
    return batches

def execute_translation_batches(target_items, model_name, progress_bar, status_text, preview_box):
    batches = dynamic_chunking(target_items)
    processed = 0
    total_targets = len(target_items)
    
    for chunk in batches:
        texts_to_translate = []
        for item in chunk:
            if smart_resume and not re.search(r'[\u4e00-\u9fff]', item['text']):
                item['translated'] = item['text'] 
            else:
                texts_to_translate.append(item['text'])
                item['translated'] = None
                
        if texts_to_translate:
            translated_results = translate_batch_core(texts_to_translate, model_name)
            res_idx = 0
            for item in chunk:
                if item['translated'] is None:
                    item['translated'] = translated_results[res_idx] if res_idx < len(translated_results) else item['text']
                    res_idx += 1
                    
        for item in chunk:
            original = item['text']
            translated = item['translated']
            if 'apply_func' in item:
                item['apply_func'](item['obj'], translated)
                
            processed += 1
            progress_bar.progress(processed / total_targets)
            preview_box.success(f"**Đã xử lý {processed}/{total_targets} mục tiêu:**\n\n🇨🇳 `{original}`\n\n🇻🇳 `{translated}`")
            
            if 'autosave_func' in item and processed % 5 == 0:
                item['autosave_func']()

# ==========================================
# CÁC HÀM XỬ LÝ ĐỊNH DẠNG FILE
# ==========================================

# 1. FILE WORD (.docx)
def process_word(file_bytes, model_name, progress_bar, status_text, preview_box, file_name):
    doc = docx.Document(io.BytesIO(file_bytes))
    autosave_filename = f"autosave_{file_name}"
    
    for section in doc.sections:
        section.page_width, section.page_height = Mm(210), Mm(297)
        section.top_margin = section.bottom_margin = section.right_margin = Cm(2.0)
        section.left_margin = Cm(3.0)

    def apply_word(p, translated_text):
        p.clear() 
        p.add_run(translated_text)
        p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        for run in p.runs:
            run.font.name = 'Times New Roman'
            try:
                rFonts = run._r.get_or_add_rPr().get_or_add_rFonts()
                rFonts.set(docx.oxml.ns.qn('w:ascii'), 'Times New Roman')
                rFonts.set(docx.oxml.ns.qn('w:hAnsi'), 'Times New Roman')
            except: pass
        try: p.paragraph_format.first_line_indent = Cm(1.27)
        except: pass

    def autosave_word():
        try: doc.save(autosave_filename)
        except: pass

    all_items = []
    seen_paras = set()
    for p_xml in doc.element.xpath('//w:p'):
        p = Paragraph(p_xml, doc)
        text = p.text.strip()
        if text and p not in seen_paras:
            is_in_shape = len(p_xml.xpath('ancestor::w:txbxContent')) > 0
            if skip_shapes and is_in_shape: continue
            all_items.append({
                'obj': p, 'text': text, 
                'apply_func': apply_word, 'autosave_func': autosave_word
            })
            seen_paras.add(p)
            
    actual_start = max(0, start_idx - 1) if dich_theo_vung else 0
    actual_end = end_idx if (dich_theo_vung and end_idx > 0) else len(all_items)
    target_items = all_items[actual_start:actual_end]
    
    if target_items:
        status_text.info(f"📊 Tổng: {len(all_items)} đoạn. Đang dịch: {len(target_items)} đoạn.")
        execute_translation_batches(target_items, model_name, progress_bar, status_text, preview_box)
        
    output = io.BytesIO()
    doc.save(output)
    if os.path.exists(autosave_filename):
        try: os.remove(autosave_filename) 
        except: pass
    return output.getvalue()

# 2. FILE EXCEL (.xlsx) - SỬA LỖI I/O
def process_excel(file_bytes, model_name, progress_bar, status_text, preview_box, file_name):
    # 1. GHI FILE TẠM RA Ổ CỨNG VẬT LÝ (Trị dứt điểm lỗi I/O RAM)
    stamp = int(time.time() * 1000)
    temp_input = f"temp_in_{stamp}.xlsx"
    temp_output = f"temp_out_{stamp}.xlsx"
    autosave_filename = f"autosave_{file_name}"
    
    with open(temp_input, "wb") as f:
        f.write(file_bytes)
        
    wb = openpyxl.load_workbook(temp_input)

    def apply_excel(cell, translated_text):
        cell.value = translated_text

    def autosave_excel():
        try: wb.save(autosave_filename)
        except: pass

    all_items = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                val = cell.value
                if isinstance(val, str) and not val.startswith('=') and val.strip():
                    all_items.append({
                        'obj': cell, 'text': val.strip(), 
                        'apply_func': apply_excel, 'autosave_func': autosave_excel
                    })

    actual_start = max(0, start_idx - 1) if dich_theo_vung else 0
    actual_end = end_idx if (dich_theo_vung and end_idx > 0) else len(all_items)
    target_items = all_items[actual_start:actual_end]
    
    if target_items:
        status_text.info(f"📊 Tổng: {len(all_items)} ô. Đang dịch: {len(target_items)} ô.")
        execute_translation_batches(target_items, model_name, progress_bar, status_text, preview_box)

    # 2. LƯU KẾT QUẢ RA FILE VẬT LÝ (Tránh lỗi đóng luồng)
    wb.save(temp_output)
    wb.close()
    
    # 3. ĐỌC LẠI FILE VẬT LÝ THÀNH BYTES ĐỂ TẢI VỀ
    with open(temp_output, "rb") as f:
        result_bytes = f.read()
        
    # 4. DỌN SẠCH RÁC TRÊN Ổ CỨNG
    try: os.remove(temp_input)
    except: pass
    try: os.remove(temp_output)
    except: pass
    if os.path.exists(autosave_filename):
        try: os.remove(autosave_filename) 
        except: pass
        
    return result_bytes

# 3. FILE POWERPOINT (.pptx)
def process_pptx(file_bytes, model_name, progress_bar, status_text, preview_box, file_name):
    prs = Presentation(io.BytesIO(file_bytes))
    autosave_filename = f"autosave_{file_name}"

    def apply_pptx(paragraph, translated_text):
        paragraph.text = translated_text

    def autosave_pptx():
        try: prs.save(autosave_filename)
        except: pass

    all_items = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame: continue
            for p in shape.text_frame.paragraphs:
                text = p.text.strip()
                if text:
                    all_items.append({
                        'obj': p, 'text': text,
                        'apply_func': apply_pptx, 'autosave_func': autosave_pptx
                    })

    actual_start = max(0, start_idx - 1) if dich_theo_vung else 0
    actual_end = end_idx if (dich_theo_vung and end_idx > 0) else len(all_items)
    target_items = all_items[actual_start:actual_end]
    
    if target_items:
        status_text.info(f"📊 Tổng: {len(all_items)} đoạn PPT. Đang dịch: {len(target_items)} đoạn.")
        execute_translation_batches(target_items, model_name, progress_bar, status_text, preview_box)

    output = io.BytesIO()
    prs.save(output)
    if os.path.exists(autosave_filename):
        try: os.remove(autosave_filename) 
        except: pass
    return output.getvalue()

# 4. FILE PDF (.pdf) -> XUẤT RA DOCX GIỮ NGUYÊN BẢNG BIỂU & HÌNH ẢNH
def process_pdf(file_bytes, model_name, progress_bar, status_text, preview_box, file_name):
    temp_pdf = f"temp_in_{file_name}"
    temp_docx = f"temp_out_{file_name.replace('.pdf', '.docx')}"
    
    with open(temp_pdf, "wb") as f:
        f.write(file_bytes)
        
    status_text.info("🔄 Đang chuyển đổi cấu trúc PDF sang Word để giữ nguyên Bảng biểu và Hình ảnh...")
    try:
        cv = Converter(temp_pdf)
        cv.convert(temp_docx, start=0, end=None)
        cv.close()
    except Exception as e:
        status_text.error(f"Lỗi chuyển đổi PDF: {e}")
        return None

    status_text.info("✅ Chuyển đổi thành công! Bắt đầu tiến hành dịch thuật...")
    
    with open(temp_docx, "rb") as f:
        docx_bytes = f.read()
        
    try: os.remove(temp_pdf)
    except: pass
    try: os.remove(temp_docx)
    except: pass
    
    # Ép luồng Word xử lý file DOCX vừa sinh ra (để tận dụng toàn bộ sức mạnh tối ưu lề, indent)
    return process_word(docx_bytes, model_name, progress_bar, status_text, preview_box, file_name.replace('.pdf', '.docx'))

# ==========================================
# CHIA TAB GIAO DIỆN CHÍNH
# ==========================================
tab1, tab2 = st.tabs(["📄 Dịch Đa Định Dạng (Word, Excel, PPTX, PDF)", "🖼️ Dịch Ảnh Định Vị"])

with tab1:
    st.header("📂 Hệ Thống Xử Lý File Kỹ Thuật")
    
    autosave_files = [f for f in os.listdir('.') if f.startswith('autosave_')]
    if autosave_files:
        st.error("⚠️ **PHÁT HIỆN FILE DỊCH DỞ DANG!** Tải file nháp về, thả vào ô bên dưới để DỊCH TIẾP:")
        for f in autosave_files:
            with open(f, "rb") as file_data:
                st.download_button(label=f"📥 TẢI XUỐNG BẢN NHÁP", data=file_data, file_name=f, mime="application/octet-stream", type="primary")
        if st.button("🗑️ Đã tải xong! Xóa bản nháp này"):
            for f in autosave_files:
                try: os.remove(f)
                except: pass
            st.rerun()
            
    uploaded_file = st.file_uploader("📂 Tải file (.docx, .xlsx, .pptx, .pdf)", type=['docx', 'xlsx', 'pptx', 'pdf'])
    
    if uploaded_file:
        file_bytes = uploaded_file.read()
        file_ext = os.path.splitext(uploaded_file.name)[1].lower()
        
        with st.expander("🔍 Bấm vào đây để SOI VỊ TRÍ văn bản trong file"):
            if st.button("Hiển thị danh sách đánh số"):
                preview_list = []
                if file_ext == '.docx':
                    temp_doc = docx.Document(io.BytesIO(file_bytes))
                    for p_xml in temp_doc.element.xpath('//w:p'):
                        p = Paragraph(p_xml, temp_doc)
                        if p.text.strip(): preview_list.append(p.text.strip())
                elif file_ext == '.xlsx':
                    temp_wb = openpyxl.load_workbook(io.BytesIO(file_bytes))
                    for sheet in temp_wb.worksheets:
                        for row in sheet.iter_rows():
                            for cell in row:
                                if isinstance(cell.value, str) and cell.value.strip(): preview_list.append(cell.value.strip())
                elif file_ext == '.pptx':
                    temp_prs = Presentation(io.BytesIO(file_bytes))
                    for slide in temp_prs.slides:
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                for p in shape.text_frame.paragraphs:
                                    if p.text.strip(): preview_list.append(p.text.strip())
                elif file_ext == '.pdf':
                    # Do PDF giờ dùng pdf2docx nên không soi trực tiếp raw text nữa
                    st.warning("⚠️ File PDF cần được chuyển đổi sang Word trước khi dịch, vì vậy hãy bỏ qua bước soi X-Quang này và bấm dịch luôn nhé.")
                
                if preview_list:
                    df = pd.DataFrame({"Số Thứ Tự": range(1, len(preview_list) + 1), "Nội dung": [t[:100] + "..." if len(t) > 100 else t for t in preview_list]})
                    st.dataframe(df, use_container_width=True, hide_index=True)

        if st.button("🚀 Bắt Đầu Dịch File", use_container_width=True):
            progress_bar = st.progress(0)
            status_text = st.empty()
            preview_box = st.empty()
            
            try:
                if file_ext == '.docx':
                    result_bytes = process_word(file_bytes, selected_model_name, progress_bar, status_text, preview_box, uploaded_file.name)
                    mime_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                    final_name = f"Translated_{uploaded_file.name}"
                elif file_ext == '.xlsx':
                    result_bytes = process_excel(file_bytes, selected_model_name, progress_bar, status_text, preview_box, uploaded_file.name)
                    mime_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                    final_name = f"Translated_{uploaded_file.name}"
                elif file_ext == '.pptx':
                    result_bytes = process_pptx(file_bytes, selected_model_name, progress_bar, status_text, preview_box, uploaded_file.name)
                    mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    final_name = f"Translated_{uploaded_file.name}"
                elif file_ext == '.pdf':
                    result_bytes = process_pdf(file_bytes, selected_model_name, progress_bar, status_text, preview_box, uploaded_file.name)
                    mime_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                    final_name = f"Translated_{uploaded_file.name.replace('.pdf', '.docx')}"
                    
                status_text.success("✅ Đã xử lý xong vùng mục tiêu!")
                preview_box.empty()
                st.download_button("📥 TẢI XUỐNG TỆP HOÀN CHỈNH", data=result_bytes, file_name=final_name, mime=mime_type, type="primary", use_container_width=True)
            except Exception as e:
                st.error(f"❌ Lỗi xử lý hệ thống: {str(e)}")

with tab2:
    st.header("🖼️ Đọc & Định Vị Chữ Trong Ảnh")
    img_to_process = None
    col1, col2 = st.columns(2)
    with col1:
        paste_result = paste_image_button(label="📋 BẤM VÀO ĐÂY ĐỂ DÁN ẢNH (CTRL+V)", text_color="#ffffff", background_color="#0068c9")
        if paste_result.image_data is not None: img_to_process = paste_result.image_data.convert('RGB')
    with col2:
        uploaded_image = st.file_uploader("Hoặc Tải File Lên", type=['jpg', 'jpeg', 'png'], label_visibility="collapsed")
        if uploaded_image is not None: img_to_process = Image.open(uploaded_image).convert('RGB')

    if img_to_process is not None:
        st.image(img_to_process, use_column_width=True)
        if st.button("🔍 TIẾN HÀNH DỊCH & ĐỊNH VỊ", type="primary", use_container_width=True):
            with st.spinner("AI đang quét ảnh..."):
                model = genai.GenerativeModel(selected_model_name)
                res = model.generate_content(["Trích xuất và dịch chữ tiếng Trung trong ảnh. Trình bày dạng danh sách: - [Vị trí tương đối] Tiếng Trung -> Tiếng Việt", img_to_process], safety_settings=SAFETY_SETTINGS)
                st.code(res.text, language="markdown")